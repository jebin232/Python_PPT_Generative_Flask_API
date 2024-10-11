from flask import Flask, request, send_file, render_template
from pptx import Presentation
import io
import google.generativeai as genai

app = Flask(__name__)

# Google Gemini API Key Configuration
gemini_api_key = "AIzaSyBct_7JTaia75avYy1TXMfTb_bfueU1vkw"
genai.configure(api_key=gemini_api_key)

def generate_text_from_title(title):
    model = genai.GenerativeModel('gemini-pro')
    
    response = model.generate_content(
        f"{title}",
        generation_config=genai.types.GenerationConfig(
            candidate_count=1,
            stop_sequences=['space'],
            max_output_tokens=1000,
            temperature=0.7
        )
    )
    
    return response.text

def create_presentation(title, contents):
    prs = Presentation()
    content = contents.replace('*', '')

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Split content into paragraphs
    paragraphs = content.split('\n')

    # Add a new slide for each paragraph
    for paragraph in paragraphs:
        if paragraph.strip():  # Only add slide for non-empty paragraphs
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            # title_placeholder.text = 'Content'
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = paragraph

    return prs

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        title = request.form['title']
        content = generate_text_from_title(title)
        prs = create_presentation(title, content)
        
        # Save presentation to a BytesIO stream
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        
        return send_file(
                    ppt_io, 
                    download_name='presentation.pptx', 
                    as_attachment=True, 
                    mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
        
    return render_template('index.html')

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8000)
